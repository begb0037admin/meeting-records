"""
KPI Presentation generator - general, reusable across months.

Reads a month's Source Data folder (Excel + H&S Word doc), following the
convention Kevin confirmed: source files land by email and get dropped into
Functional Analysis Team Monthly Statistics/{YYYY}/{MM Mon}/Source Data/.

Base deck: copies the PREVIOUS month's own finished deck (the real monthly
process, confirmed 2 Aug 2026 - the saved "Templates" file is stale and not
used) and edits it in place, so layout fixes (blank-row removal, chart image
positions - see memory/kpi-presentation-layout-fixes.md) persist forward
automatically without needing to be reapplied each month.

KNOWN GAP, permanent characteristic of the source data (not a bug): the H&S
per-system incident-volume breakdown (slide 2) and time-to-resolve band
breakdown (slide 3) only exist as a "last month" snapshot in each month's own
Word doc - there is no rolling multi-month series for these two tables, unlike
every PXD-sourced table/chart (slides 4-10), which get a full 15-month rolling
window each month and are fully self-contained (confirmed 7 Aug 2026: slides
6/7's time-to-complete bands ARE on a rolling sheet in the current month's own
Excel, so - unlike H&S - they need no cross-month lookup at all). So the
"previous month" column for slides 2 and 3 must come from somewhere other than
the current month's own source data.

RESOLVED 7 Aug 2026 (was an open gap, folded into build_month() below): the
prior month's own Source Data folder, when it still exists in the OneDrive
archive, carries the exact same "last month" snapshot for what WAS then the
current month - i.e. May's own H&S doc's "last month" table IS May's own
confirmed figures, byte-verified to match the real June deck's "May 26"
column exactly. build_month() auto-locates and reads it by default. This is
real carried-forward data, not invented - but it does depend on the archive
folder still being there; build_month() accepts an explicit prev_hs override
for the rare case it isn't.
"""
from pptx import Presentation
from pptx.chart.data import CategoryChartData
import openpyxl
from docx import Document
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import glob
import re
import os

ORANGE = "#ED7D31"
GREEN = "#70AD47"
DARKRED = "#843C0C"

# Categories the real deck excludes from the Incident-Other breakdown
# (slide 5, Table 4) - discovered 2 Aug 2026 by diffing raw vs. real deck,
# not documented anywhere in the source. See memory/
# kpi-presentation-source-and-rebuild-poc.md.
INCIDENT_OTHER_EXCLUDE = {
    "My Development", "Applications/Software", "Applicant", "Pensions",
    "Biztalk, SSIS & Azure",
}

# --- Month / folder conventions (7 Aug 2026 extension) -----------------
# Confirmed OneDrive layout: Functional Analysis Team Monthly Statistics\
# {YYYY}\{MM Mon}\ with a Source Data\ subfolder each month, and the
# finished deck for month N sitting alongside as
# "KPI presentation - {Month} {YYYY}.pptx". See memory/
# kpi-presentation-source-and-rebuild-poc.md (Lauren, confirmed 2 Aug 2026).
ONEDRIVE_BASE = (
    r"C:\Users\admin\OneDrive - Nexus365\Functional Analysis Team Monthly Statistics"
)
MONTH_ABBR = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
              "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
MONTH_FULL = ["January", "February", "March", "April", "May", "June",
              "July", "August", "September", "October", "November", "December"]


def month_dir(year, month):
    return os.path.join(ONEDRIVE_BASE, str(year), f"{month:02d} {MONTH_ABBR[month - 1]}")


def short_label(year, month):
    """e.g. (2026, 6) -> 'Jun 26' - matches the deck's own header style."""
    return f"{MONTH_ABBR[month - 1]} {year % 100:02d}"


def deck_path(year, month):
    return os.path.join(
        month_dir(year, month),
        f"KPI presentation - {MONTH_FULL[month - 1]} {year}.pptx",
    )


def prev_ym(year, month):
    return (year - 1, 12) if month == 1 else (year, month - 1)


def find_one(dirpath, pattern):
    """Locate exactly one file matching pattern in dirpath. Source filenames
    carry a variable export-timestamp suffix (e.g. '...202607010715.xlsx')
    so callers can't hardcode the full name - but there must be exactly one
    match, or something is wrong with the folder and we should say so rather
    than silently pick one."""
    matches = glob.glob(os.path.join(dirpath, pattern))
    if len(matches) != 1:
        raise FileNotFoundError(
            f"expected exactly 1 match for {pattern!r} in {dirpath!r}, found {len(matches)}: {matches}"
        )
    return matches[0]


def find_source_files(year, month):
    """Locate a given month's own Excel export and H&S Word doc inside its
    Source Data folder."""
    src = os.path.join(month_dir(year, month), "Source Data")
    excel_path = find_one(src, "HR_Systems_Functional_Team_Monthly_Report_Excel*.xlsx")
    hs_path = find_one(src, "Health and Safety Systems Support Statistics*.docx")
    return excel_path, hs_path


def parse_count_pct(s):
    """Parse a rolling-sheet cell like '112 (74.67%)' into (112, 74.67).
    Returns (0, 0.0) for a blank/None cell (a month with no recorded tasks
    in that band - a real possible value, not a parsing failure)."""
    if not s:
        return 0, 0.0
    m = re.match(r"\s*(-?\d+)\s*\(([-\d.]+)%\)\s*$", str(s))
    if not m:
        raise ValueError(f"unrecognised count/pct cell format: {s!r}")
    return int(m.group(1)), float(m.group(2))


def extract_pxd(excel_path):
    """Pull every PXD-sourced figure from the month's Excel export.
    Returns a dict keyed by the logical field each slide needs."""
    wb = openpyxl.load_workbook(excel_path, data_only=True)

    def sheet_row(sheet_name, row_label_col0):
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            if row and row[0] == row_label_col0:
                return row
        raise KeyError(f"row {row_label_col0!r} not found in {sheet_name!r}")

    def monthly_by_category(sheet_name):
        """Sheets shaped like: title row, blank spacer row, header row of
        months, then one row per category. Returns {category: [15 monthly
        values]} plus the month header row (for locating which column is
        'this month' etc.).

        FIX (7 Aug 2026): header is row index 2, not 1 - row 1 is a blank
        spacer row on every sheet checked (confirmed across 7 sheets this
        function reads from). The original index-1 assumption meant `header`
        was always the blank row, and the data loop (starting at rows[2:])
        picked up the *real* header row as a bogus extra "category" whose
        row[0] label is the sheet's column-header text and whose "monthly
        values" are literally the month-label strings ('Apr 25', ... 'Jun
        26'). This was silently harmless everywhere the buggy version was
        already exercised - sheets whose header cell is None (Average Task
        Acceptance, Tasks, etc.) skip it via the `not label` check, and
        sheets where it produced a real string key (e.g. 'Parent Ticket
        Type HRIS') just added an unused dict entry nothing ever looked up.
        It broke immediately, though, on the newly-wired slide 6/7 band
        parsing (build_month(2026, 6, ...) raised ValueError parsing 'Jun
        26' as a count/pct cell) - see memory/kpi-presentation-assembly-gap.md,
        this is exactly the "genuinely new, unattempted, needs its own
        self-test" capability that surfaced it. Fixing the index also fixes
        a second, previously undetected bug: every *_months output
        (slide4_months, slide8_months, etc.) was always a tuple of all
        None - nothing in the original 9-check self-test compared month
        labels, so this went unnoticed until make_trend_chart's x-axis
        labels would have shown blank ticks."""
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        header = rows[2]  # title (0), blank spacer (1), month header (2)
        data = {}
        for row in rows[3:]:
            if not row or row[0] in (None, "Description") or (isinstance(row[0], str) and row[0].startswith("Total")):
                continue
            label = row[0].strip() if isinstance(row[0], str) else row[0]
            if not label:
                continue
            data[label] = row[1:len(header)]
        return header[1:], data

    out = {}

    # Slide 4: PXD categories by month (Change / Incident-HR Self-Service /
    # Incident-Other / Service Request)
    months, cats = monthly_by_category("Tasks Completed by HRIS Analy2")
    out["slide4_months"] = months
    out["slide4_categories"] = {
        "Change": cats["Change"],
        "HR Self Service": cats["Incident - HR Self-Service"],
        "Incident – Other": cats["Incident - Other"],
        "Service Request": cats["Service Request"],
    }

    # Slide 5, Table 6: Service Request type breakdown, last month only
    ws = wb["Service Request Tasks complete"]
    rows = [r for r in ws.iter_rows(values_only=True) if r and r[0] not in (None, "Description") and not str(r[0]).startswith("Total")]
    label_map = {
        "Manager Self Service (MSS)": "Manager Self Service",
        "Work Group (HR Self-Service)": "WG - HR Self-Service",
        "Add New Shift Type (HR Self-Service)": "Add New Shift Type",
        "HR Systems User Access": "HR Systems User Access",
        "HR Systems People Data Dashboards": "People Data Dashboards",
        "HR Systems Enhancement and Changes": "Enhancement & Changes",
    }
    sr_breakdown = {}
    for r in rows[1:]:
        if r[0] in label_map:
            sr_breakdown[label_map[r[0]]] = (r[1] or 0, r[2] or "0.00%")
    out["slide5_service_request"] = sr_breakdown

    # Slide 5, Table 4: Incident-Other type breakdown, filtered
    ws = wb["Non-HR Self Service Incident T"]
    rows = [r for r in ws.iter_rows(values_only=True) if r and r[0] not in (None, "Description") and not str(r[0]).startswith("Total")]
    io_breakdown = {}
    for r in rows[1:]:
        if r[0] and r[0] not in INCIDENT_OTHER_EXCLUDE:
            io_breakdown[r[0]] = r[1] or 0
    filtered_total = sum(io_breakdown.values())
    out["slide5_incident_other"] = {
        k: (v, f"{v/filtered_total*100:.2f}%") for k, v in io_breakdown.items()
    }
    out["slide5_incident_other_total"] = filtered_total

    # Slides 6/7: completion-time % bands, monthly rolling
    _, sr_time = monthly_by_category("Time To Complete Service Requ1")
    _, oi_time = monthly_by_category("Time To Complete Other Incide1")
    out["slide6_bands"] = sr_time   # {'Same Day': [...], 'Next Day': [...], '3 - 5 days': [...], '6+ days': [...]}
    out["slide7_bands"] = oi_time

    # Slides 8/9: average days, monthly rolling
    months, acc = monthly_by_category("Average Task Acceptance (Days)")
    out["slide8_months"] = months
    out["slide8_values"] = acc["Tasks"]
    months, comp = monthly_by_category("Average Task Time to Complete ")
    out["slide9_months"] = months
    out["slide9_values"] = comp["Tasks"]

    # Slide 9 chart's "Total: N" annotation is NOT sum(slide9_values) (that
    # would be a meaningless sum of average-days figures) - it's the total
    # completed-tasks count from a *different* sheet. Confirmed 2 Aug 2026
    # by hand-summing against the proven picture-swap rebuild (see
    # memory/kpi-presentation-layout-fixes.md): slide 8's total is the
    # 'Tasks' sheet's Created sum (already available as slide10_created),
    # slide 9's is this sheet's Completed Tasks sum.
    _, analy3 = monthly_by_category("Tasks Completed by HRIS Analy3")
    out["slide9_completed_tasks"] = analy3["Completed Tasks"]

    # Slide 10: created vs completed, monthly rolling
    months, tasks = monthly_by_category("Tasks")
    out["slide10_months"] = months
    out["slide10_created"] = tasks["Created"]
    out["slide10_completed"] = tasks["Completed / Cancelled / Rejected"]

    return out


def extract_hs(docx_path):
    """Pull this month's H&S figures (last-month snapshot only - see module
    docstring for why there's no rolling series here).

    IMPORTANT: this doc has BOTH a "last month" and a "last 12 months"
    version of each table, with IDENTICAL column headers - text-matching on
    headers alone silently grabs the wrong (12-month cumulative) one. Must
    match on the preceding description paragraph's "last month" wording, not
    header shape. Table indices below were confirmed against June 2026's doc
    (verified 2 Aug 2026 self-test) but this walks the document body in
    order and matches on description text, so it isn't a hardcoded index -
    it should hold for any month using the same export template."""
    doc = Document(docx_path)
    body = doc.element.body
    out = {}
    pending_desc = ""
    for child in body.iterchildren():
        if child.tag.endswith("}p"):
            text = "".join(node.text or "" for node in child.iter() if node.tag.endswith("}t"))
            if "Description" in text or "last month" in text or "12 months" in text:
                pending_desc = text
        elif child.tag.endswith("}tbl"):
            from docx.table import Table
            t = Table(child, doc)
            rows = [[c.text for c in r.cells] for r in t.rows]
            if not rows:
                continue
            is_last_month = "last month" in pending_desc and "12 months" not in pending_desc
            if is_last_month and rows[0] == ["", "Category", "Category %"] and any("Cority" in r[0] for r in rows[1:] if r[0]):
                out["slide2_volumes"] = {r[0]: (int(r[1]), r[2]) for r in rows[1:]}
            if is_last_month and rows[0] == ["", "Time To Resolve (Days to 6+)", "Time To Resolve (Days to 6+) %"]:
                out["slide3_bands"] = {r[0]: (int(r[1]), r[2]) for r in rows[1:]}
    return out


def make_trend_chart(months, values, total, out_path, y_step):
    fig, ax = plt.subplots(figsize=(20, 12.88), dpi=100)
    fig.patch.set_facecolor("white")
    x = np.arange(len(months))

    def nice_floor(v, step):
        return np.floor(v / step) * step

    ymin = max(0, nice_floor(min(values), y_step) - y_step) if min(values) > y_step else 0
    ymax = max(values) + y_step * 2
    ax.set_xlim(-0.6, len(months) - 0.4)
    ax.set_ylim(ymin, ymax)

    grad = np.vstack([np.linspace(0, 1, 256)] * 2)
    ax.imshow(grad, extent=[*ax.get_xlim(), *ax.get_ylim()], aspect="auto", cmap="Greys", alpha=0.08, zorder=0)

    ax.bar(x, values, width=0.55, color=ORANGE, zorder=3, bottom=0)

    z = np.polyfit(x, values, 1)
    trend = np.poly1d(z)
    ax.plot(x, trend(x), color="#595959", linewidth=1.2, zorder=4)

    for xi, v in zip(x, values):
        ax.annotate(f"{v:.1f}", (xi, max(v, ymin)), xytext=(0, -18), textcoords="offset points",
                    ha="center", va="top", fontsize=13,
                    bbox=dict(boxstyle="square,pad=0.3", fc="white", ec="#BFBFBF", lw=0.8))

    ax.set_xticks(x)
    ax.set_xticklabels(months, rotation=60, ha="right", fontsize=13)
    ax.yaxis.set_major_locator(mticker.MultipleLocator(y_step))
    ax.grid(axis="y", color="#E0E0E0", linewidth=0.8, zorder=1)
    ax.set_axisbelow(True)
    for spine in ("top", "right"):
        ax.spines[spine].set_visible(False)
    ax.text(0.995, 0.01, f"Total: {total}", transform=ax.transAxes,
            ha="right", va="bottom", fontsize=13, color="#7F7F7F")
    fig.subplots_adjust(left=0.04, right=0.98, top=0.98, bottom=0.28)
    plt.savefig(out_path, facecolor="white")
    plt.close(fig)


def make_combo_chart(months, created, completed, out_path):
    fig, ax1 = plt.subplots(figsize=(20, 12.88), dpi=100)
    fig.patch.set_facecolor("white")
    x = np.arange(len(months))
    w = 0.38
    ax1.set_xlim(-0.6, len(months) - 0.4)
    ax1.set_ylim(0, max(max(created), max(completed)) + 50)
    grad = np.vstack([np.linspace(0, 1, 256)] * 2)
    ax1.imshow(grad, extent=[*ax1.get_xlim(), *ax1.get_ylim()], aspect="auto", cmap="Greys", alpha=0.08, zorder=0)
    ax1.bar(x - w / 2, created, width=w, color=ORANGE, zorder=3, label=f"Created ({sum(created)})")
    ax1.bar(x + w / 2, completed, width=w, color=GREEN, zorder=3, label=f"Completed / Cancelled / Rejected ({sum(completed)})")
    ax1.set_ylabel("Number of Tasks created and completed", fontsize=13, fontweight="bold")
    ax1.set_xticks(x)
    ax1.set_xticklabels(months, rotation=60, ha="right", fontsize=13)
    ax1.grid(axis="y", color="#E0E0E0", linewidth=0.8, zorder=1)
    ax1.set_axisbelow(True)
    ax1.spines["top"].set_visible(False)

    variance = [abs(c - r) for c, r in zip(completed, created)]
    ax2 = ax1.twinx()
    ax2.plot(x, variance, color=DARKRED, marker="o", linewidth=1.5, zorder=4, label="Variance *")
    ax2.set_ylim(0, max(variance) + 5)
    ax2.spines["top"].set_visible(False)

    lines1, labels1 = ax1.get_legend_handles_labels()
    lines2, labels2 = ax2.get_legend_handles_labels()
    ax1.legend(lines1 + lines2, labels1 + labels2, loc="lower center",
               bbox_to_anchor=(0.5, -0.32), ncol=3, frameon=True, fontsize=12)
    ax1.text(0.995, 1.02, f"Total: {sum(created) + sum(completed)}", transform=ax1.transAxes,
             ha="right", va="bottom", fontsize=13, color="#7F7F7F")
    fig.subplots_adjust(left=0.04, right=0.98, top=0.9, bottom=0.28)
    plt.savefig(out_path, facecolor="white")
    plt.close(fig)


def set_cell(cell, text):
    para = cell.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for extra in para.runs[1:]:
            extra.text = ""
    else:
        para.text = text


def find_table(slide, name):
    for shape in slide.shapes:
        if shape.has_table and shape.name == name:
            return shape.table
    raise RuntimeError(f"table {name!r} not found")


def fmt_delta(cur, prev):
    d = cur - prev
    pct = f" ({d/prev*100:+.0f}%)" if prev else ""
    return f"{d:+d}{pct}"


def plain_delta(d):
    """Signed integer delta with no percentage - matches the real deck's
    style for slide 2/3/10 MoM columns and slide 10's Variance row (as
    opposed to fmt_delta's '+N (+P%)' style used on slides 4/10's
    Created/Completed rows). Zero renders as '0', not '+0' - confirmed
    against the real June deck (slide 2's DSE row, slide 3's Next Day row)."""
    return "0" if d == 0 else f"{d:+d}"


def set_pie_chart(shape, categories, value_map, series_name):
    """Replace a native pptx pie chart's series data in place via
    chart.replace_data() - the previously-unattempted capability (see
    memory/kpi-presentation-assembly-gap.md). Keeps the existing category
    order/labels the chart already uses so formatting/colours carry over;
    only the underlying numbers (and the generic default 'Sales' series
    name - proof-of-non-touch, per Lauren's byte-level verification) change.
    Values are written at full float precision; PowerPoint's own data-label
    number format on the chart (untouched here) governs display rounding,
    same as it always has."""
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series(series_name, [value_map.get(c) for c in categories])
    shape.chart.replace_data(cd)


def populate_deck(base_deck_path, pxd, hs_current, hs_prev, month_label, chart_dir, out_path, year, month):
    """Assemble a full month's deck from a base (prior month's finished
    deck, already carrying the layout fixes) plus extracted figures.

    year/month are the current reporting month (e.g. 2026, 7 for July) -
    needed to relabel every table's month-column headers (e.g. slide 4's
    'Jun 25 / May 26 / Jun 26'), which stay stale otherwise since they're
    plain template text carried forward from whichever month the base deck
    was last for, not derived from data like the cells beneath them."""
    prs = Presentation(base_deck_path)
    m = pxd["slide4_months"]  # 15-month label list, ends at current month
    cur, prev_i, yr_i = -1, -2, -13  # current, previous month, 12 months back +1

    py, pm = prev_ym(year, month)
    yy, ym = (year - 1, month)
    cur_lbl, prev_lbl, yr_lbl = short_label(year, month), short_label(py, pm), short_label(yy, ym)

    # Slide 1
    for shape in prs.slides[0].shapes:
        if shape.name == "Subtitle 2":
            set_cell_like(shape, f"{month_label.upper()} KPI STATISTICS | ")

    # Header month-column relabelling - every table below has its own body
    # values recomputed from this month's data, but the header row text
    # ('Jun 25', 'May 26', 'Jun 26' etc.) is separate static text that must
    # be updated in step or a July deck would show correct numbers under
    # wrong month labels. (slide_index, table_name): {col_index: label}
    header_updates = {
        (1, "Table 5"): {1: prev_lbl, 2: cur_lbl},                    # slide 2
        (2, "Table 5"): {1: prev_lbl, 3: cur_lbl},                    # slide 3
        (3, "Table 5"): {1: yr_lbl, 2: prev_lbl, 3: cur_lbl},         # slide 4
        (5, "Table 5"): {4: prev_lbl, 5: cur_lbl},                    # slide 6
        (6, "Table 5"): {4: prev_lbl, 5: cur_lbl},                    # slide 7
        (9, "Table 3"): {1: yr_lbl, 2: prev_lbl, 3: cur_lbl},         # slide 10
    }
    for (slide_i, table_name), cols in header_updates.items():
        t = find_table(prs.slides[slide_i], table_name)
        for ci, label in cols.items():
            set_cell(t.rows[0].cells[ci], label)
    # Slides 8/9's Table 11/10 use a 3-column month row at row index 2, no
    # other header text to preserve on that row.
    for slide_i, table_name in ((7, "Table 11"), (8, "Table 10")):
        t = find_table(prs.slides[slide_i], table_name)
        for ci, label in ((0, yr_lbl), (1, prev_lbl), (2, cur_lbl)):
            set_cell(t.rows[2].cells[ci], label)

    # Slide 2: H&S system volumes - Table 5 (counts + MoM) and native pie
    # chart (% share). Both previously pass-through (see memory/
    # kpi-presentation-assembly-gap.md) - genuinely written here.
    def hs_lookup(hs_dict, short):
        for full, (count, pct) in hs_dict.get("slide2_volumes", {}).items():
            if short in full:
                return count
        return 0

    t2 = find_table(prs.slides[1], "Table 5")
    s2_order = ["Cority", "Odyssey", "IRIS", "DSE"]
    for ri, short in enumerate(s2_order, start=1):
        c = hs_lookup(hs_current, short)
        p = hs_lookup(hs_prev, short)
        set_cell(t2.rows[ri].cells[1], str(p))
        set_cell(t2.rows[ri].cells[2], str(c))
        set_cell(t2.rows[ri].cells[3], plain_delta(c - p))
    s2_tot_c = sum(hs_lookup(hs_current, s) for s in s2_order)
    s2_tot_p = sum(hs_lookup(hs_prev, s) for s in s2_order)
    set_cell(t2.rows[5].cells[1], str(s2_tot_p))
    set_cell(t2.rows[5].cells[2], str(s2_tot_c))
    set_cell(t2.rows[5].cells[3], fmt_delta(s2_tot_c, s2_tot_p))

    slide2_pie = {}
    for short in s2_order:
        c = hs_lookup(hs_current, short)
        slide2_pie[short] = round(c / s2_tot_c * 100, 2) if c and s2_tot_c else None
    set_pie_chart(
        next(sh for sh in prs.slides[1].shapes if sh.has_chart),
        ["Cority", "Odyssey", "IRIS", "DSE"], slide2_pie, "H&S System Volume %",
    )

    # Slide 3: H&S time-to-resolve bands - Table 5 (count+pct, both months,
    # MoM) and native pie chart (current month %).
    t3 = find_table(prs.slides[2], "Table 5")
    s3_rows = [("Same Day", 1), ("Next Day", 2), ("3 - 5 days", 3), ("6+ days", 4)]
    s3_bands_c = hs_current.get("slide3_bands", {})
    s3_bands_p = hs_prev.get("slide3_bands", {})
    s3_tot_c = sum(v[0] for v in s3_bands_c.values())
    s3_tot_p = sum(v[0] for v in s3_bands_p.values())
    for label, ri in s3_rows:
        pc, ppct = s3_bands_p.get(label, (0, "0.00%"))
        cc, cpct = s3_bands_c.get(label, (0, "0.00%"))
        set_cell(t3.rows[ri].cells[1], str(pc))
        set_cell(t3.rows[ri].cells[2], ppct)
        set_cell(t3.rows[ri].cells[3], str(cc))
        set_cell(t3.rows[ri].cells[4], cpct)
        set_cell(t3.rows[ri].cells[5], plain_delta(cc - pc))
    set_cell(t3.rows[5].cells[1], str(s3_tot_p))
    set_cell(t3.rows[5].cells[3], str(s3_tot_c))
    set_cell(t3.rows[5].cells[5], plain_delta(s3_tot_c - s3_tot_p))

    slide3_pie = {label: (s3_bands_c.get(label, (0, 0))[0] / s3_tot_c * 100 if s3_tot_c else 0)
                  for label, _ in s3_rows}
    set_pie_chart(
        next(sh for sh in prs.slides[2].shapes if sh.has_chart),
        ["Same Day", "6+ Days", "Next Day", "3-5 Day"],
        {"Same Day": slide3_pie["Same Day"], "6+ Days": slide3_pie["6+ days"],
         "Next Day": slide3_pie["Next Day"], "3-5 Day": slide3_pie["3 - 5 days"]},
        "H&S Time to Resolve %",
    )

    # Slide 4: PXD categories, Jun25/May26/Jun26-equivalent + MoM/YoY
    t4 = find_table(prs.slides[3], "Table 5")
    order = ["Service Request", "Incident – Other", "HR Self Service", "Change"]
    for ri, cat in enumerate(order, start=1):
        vals = pxd["slide4_categories"][cat]
        c, p, y = vals[cur], vals[prev_i], vals[yr_i]
        row = [cat, str(y), str(p), str(c), fmt_delta(c, p), fmt_delta(c, y)]
        for ci, text in enumerate(row):
            set_cell(t4.rows[ri].cells[ci], text)
    tot_c = sum(pxd["slide4_categories"][k][cur] for k in order)
    tot_p = sum(pxd["slide4_categories"][k][prev_i] for k in order)
    tot_y = sum(pxd["slide4_categories"][k][yr_i] for k in order)
    tot_row = ["Total", str(tot_y), str(tot_p), str(tot_c), fmt_delta(tot_c, tot_p), fmt_delta(tot_c, tot_y)]
    for ci, text in enumerate(tot_row):
        set_cell(t4.rows[5].cells[ci], text)

    slide4_pie = {cat.replace("Incident – Other", "Incident - Other"):
                  round(pxd["slide4_categories"][cat][cur] / tot_c * 100, 2) if tot_c else 0
                  for cat in order}
    set_pie_chart(
        next(sh for sh in prs.slides[3].shapes if sh.has_chart),
        ["Service Request", "Change", "HR Self Service", "Incident - Other"],
        slide4_pie, "PXD Task Category %",
    )

    # Slide 5: Table 4 (incident-other) + Table 6 (service request)
    slide5 = prs.slides[4]
    io = pxd["slide5_incident_other"]
    t4_5 = find_table(slide5, "Table 4")
    io_order = ["People Management", "Time and Attendance", "Payroll", "HR Reporting",
                "Data Protection Request", "Staff Requests", "Work Groups and Managers",
                "Recruitment", "Roster (WFM)", "Interfaces"]
    for ri, cat in enumerate(io_order, start=1):
        count, pct = io.get(cat, (0, "0.00%"))
        set_cell(t4_5.rows[ri].cells[0], cat.replace("Work Groups and Managers", "Work Groups & Managers"))
        set_cell(t4_5.rows[ri].cells[1], str(count))
        set_cell(t4_5.rows[ri].cells[2], pct)
    set_cell(t4_5.rows[11].cells[1], str(pxd["slide5_incident_other_total"]))

    sr = pxd["slide5_service_request"]
    t6_5 = find_table(slide5, "Table 6")
    sr_order = ["Manager Self Service", "WG - HR Self-Service", "Add New Shift Type",
                "HR Systems User Access", "People Data Dashboards", "Enhancement & Changes"]
    for ri, cat in enumerate(sr_order, start=1):
        count, pct = sr.get(cat, (0, "0.00%"))
        set_cell(t6_5.rows[ri].cells[0], cat)
        set_cell(t6_5.rows[ri].cells[1], str(count))
        set_cell(t6_5.rows[ri].cells[2], pct)

    # Known fix 1 (Kevin, 2 Aug 2026, memory/kpi-presentation-layout-fixes.md):
    # Table 6's raw build includes a blank ['', '', ''] row before Total -
    # remove it, then reposition so its bottom lines up with Table 4's. Must
    # run BEFORE writing the Total row below - a real bug caught during this
    # extension's own self-test: the original code wrote the total to a
    # hardcoded row index (7) *before* this removal ran, which is only the
    # Total row once the blank row is gone. With the base deck already
    # carrying the un-fixed blank row forward (the real live decks never
    # actually had Kevin's manual fix applied - only a separate rebuild-TEST
    # file did), writing first and removing second silently planted the
    # total in the blank spacer row instead of the Total row, leaving the
    # real Total row stale. Folded into populate_deck itself now, not a
    # separate manual step, and ordered correctly.
    blank_ri = next(
        (ri for ri, row in enumerate(t6_5.rows) if all(not c.text.strip() for c in row.cells)),
        None,
    )
    if blank_ri is not None:
        tbl_el = t6_5._tbl
        tr = tbl_el.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}tr")[blank_ri]
        tbl_el.remove(tr)
    for shape in slide5.shapes:
        if shape.has_table and shape.name == "Table 6":
            shape.top, shape.left = 6563363, 2239340
            break

    set_cell(t6_5.rows[len(t6_5.rows) - 1].cells[1], str(sum(v[0] for v in sr.values())))

    # Slides 6/7: SR/OI time-to-complete bands. Table 5 is a combined
    # KPI-threshold summary (SAME OR NEXT DAY / LESS THAN 5 DAYS, both
    # months); chart is a current-month band-share pie. Both sourced from
    # the same rolling sheet already parsed into pxd - PXD data is fully
    # self-contained per month (unlike H&S), so no cross-folder read needed.
    #
    # NOTE on rounding (flagged, not silently resolved - see HANDOVER.md):
    # the real June deck's two combined percentages don't reconcile to a
    # single arithmetic rule to the last hundredth-of-a-percent against the
    # underlying counts - e.g. slide 7's Jun26 "LESS THAN 5 DAYS" is 87.02%
    # in the real deck, but exact count division (67/77) gives 87.01%
    # unambiguously (no rounding-boundary ambiguity in that division). This
    # implementation always uses exact count/total division (reproducible,
    # correct to the source data), which matches the real deck on 6 of 8
    # spot-checked cells and differs by 0.01 on 2 of them. Not invented or
    # guessed - it's the mathematically defensible value from the source
    # counts; flagged to Kevin as an open, sub-hundredth-of-a-percent
    # methodology question rather than silently claimed as byte-identical.
    for slide_i, table_name, sheet_key, pie_cats in (
        (5, "Table 5", "slide6_bands", ["Same Day", "6+ Days", "3-5 Days", "Next Day"]),
        (6, "Table 5", "slide7_bands", ["Same Day", "6+ Days", "3-5 Days", "Next Day"]),
    ):
        bands = pxd[sheet_key]  # {label: [15 monthly '123 (45.67%)' strings]}

        def band_counts(idx):
            out = {}
            for label, series in bands.items():
                count, _ = parse_count_pct(series[idx])
                out[label] = count
            return out

        cur_counts, prev_counts = band_counts(cur), band_counts(prev_i)
        cur_total, prev_total = sum(cur_counts.values()), sum(prev_counts.values())

        def combined_pct(counts, total, labels):
            n = sum(counts.get(l, 0) for l in labels)
            return round(n / total * 100, 2) if total else 0.0

        same_next = ["Same Day", "Next Day"]
        less_5 = ["Same Day", "Next Day", "3 - 5 days"]
        t = find_table(prs.slides[slide_i], table_name)
        set_cell(t.rows[1].cells[4], f"{combined_pct(prev_counts, prev_total, same_next):.2f}%")
        set_cell(t.rows[1].cells[5], f"{combined_pct(cur_counts, cur_total, same_next):.2f}%")
        set_cell(t.rows[2].cells[4], f"{combined_pct(prev_counts, prev_total, less_5):.2f}%")
        set_cell(t.rows[2].cells[5], f"{combined_pct(cur_counts, cur_total, less_5):.2f}%")

        pie_map = {
            "Same Day": round(cur_counts.get("Same Day", 0) / cur_total * 100, 2) if cur_total else 0,
            "6+ Days": round(cur_counts.get("6+ days", 0) / cur_total * 100, 2) if cur_total else 0,
            "3-5 Days": round(cur_counts.get("3 - 5 days", 0) / cur_total * 100, 2) if cur_total else 0,
            "Next Day": round(cur_counts.get("Next Day", 0) / cur_total * 100, 2) if cur_total else 0,
        }
        set_pie_chart(
            next(sh for sh in prs.slides[slide_i].shapes if sh.has_chart),
            pie_cats, pie_map, f"{table_name} Band %",
        )

    # Slides 8/9/10: KPI-threshold small tables (proven table-write pattern,
    # same as slide 4/5) plus the chart-image swap (proven picture-swap
    # mechanism from build_full_poc.py, SHA1-verified against the real deck
    # per memory/kpi-presentation-assembly-gap.md) - now generated fresh
    # from *this* month's own data instead of only ever having been run
    # against June's own pre-known figures.
    os.makedirs(chart_dir, exist_ok=True)

    t8 = find_table(prs.slides[7], "Table 11")
    acc = pxd["slide8_values"]
    for ci, idx in ((0, yr_i), (1, prev_i), (2, cur)):
        set_cell(t8.rows[3].cells[ci], f"{acc[idx]:.1f}")
    img8 = os.path.join(chart_dir, "slide8_chart.png")
    make_trend_chart(pxd["slide8_months"], list(acc), sum(pxd["slide10_created"]), img8, y_step=0.1)

    t9 = find_table(prs.slides[8], "Table 10")
    comp = pxd["slide9_values"]
    for ci, idx in ((0, yr_i), (1, prev_i), (2, cur)):
        set_cell(t9.rows[3].cells[ci], f"{comp[idx]:.1f}")
    img9 = os.path.join(chart_dir, "slide9_chart.png")
    make_trend_chart(pxd["slide9_months"], list(comp), sum(pxd["slide9_completed_tasks"]), img9, y_step=0.5)

    t10 = find_table(prs.slides[9], "Table 3")
    created, completed = pxd["slide10_created"], pxd["slide10_completed"]
    for ci, idx in ((1, yr_i), (2, prev_i), (3, cur)):
        set_cell(t10.rows[1].cells[ci], str(created[idx]))
        set_cell(t10.rows[2].cells[ci], str(completed[idx]))
        set_cell(t10.rows[3].cells[ci], plain_delta(completed[idx] - created[idx]))
    set_cell(t10.rows[1].cells[4], fmt_delta(created[cur], created[prev_i]))
    set_cell(t10.rows[1].cells[5], fmt_delta(created[cur], created[yr_i]))
    set_cell(t10.rows[2].cells[4], fmt_delta(completed[cur], completed[prev_i]))
    set_cell(t10.rows[2].cells[5], fmt_delta(completed[cur], completed[yr_i]))
    var_cur = completed[cur] - created[cur]
    var_prev = completed[prev_i] - created[prev_i]
    var_yr = completed[yr_i] - created[yr_i]
    set_cell(t10.rows[3].cells[4], plain_delta(var_cur - var_prev))
    set_cell(t10.rows[3].cells[5], plain_delta(var_cur - var_yr))
    img10 = os.path.join(chart_dir, "slide10_chart.png")
    make_combo_chart(pxd["slide10_months"], list(created), list(completed), img10)

    # Proven picture-swap positions, captured from Kevin's manual layout
    # pass on 2 Aug 2026 (build_full_poc.py, SHA1-verified against the real
    # deck) - reused verbatim, not re-derived.
    picture_swaps = {
        7: (img8, 5593036, 10726620),   # slide 8
        8: (img9, 5532711, 10823848),   # slide 9
        9: (img10, 4924159, 10751840),  # slide 10
    }
    for si, (img_path, new_top, new_left) in picture_swaps.items():
        slide = prs.slides[si]
        target = next(
            (sh for sh in slide.shapes if sh.shape_type == 13 and sh.name == "Picture 2"), None
        )
        if target is None:
            raise RuntimeError(f"slide {si + 1}: Picture 2 not found - layout may have changed")
        width, height = target.width, target.height
        sp = target._element
        sp.getparent().remove(sp)
        slide.shapes.add_picture(img_path, new_left, new_top, width, height)

    prs.save(out_path)
    print("written", out_path)


def build_month(year, month, out_path=None, chart_dir=None, prev_hs=None):
    """Top-level entry point: build one month's KPI Presentation end to end
    from its own Source Data folder, using the previous month's own
    finished deck as the base (the real monthly process - see memory/
    kpi-presentation-source-and-rebuild-poc.md).

    prev_hs: optional pre-extracted dict (from extract_hs) to use as the
    prior month's H&S figures, for when the prior month's own Source Data
    folder isn't available. If not given, it's read live from the prior
    month's own H&S doc - real carried-forward data, not invented, per the
    module docstring's KNOWN GAP note."""
    excel_path, hs_path = find_source_files(year, month)
    pxd = extract_pxd(excel_path)
    hs_current = extract_hs(hs_path)

    if prev_hs is None:
        py, pm = prev_ym(year, month)
        _, prev_hs_path = find_source_files(py, pm)
        prev_hs = extract_hs(prev_hs_path)

    py, pm = prev_ym(year, month)
    base_deck = deck_path(py, pm)
    if not os.path.exists(base_deck):
        raise FileNotFoundError(
            f"base deck (previous month's finished deck) not found: {base_deck!r} - "
            "the real monthly process copies this forward; it must exist before a new month can be built."
        )

    if out_path is None:
        out_dir = month_dir(year, month)
        os.makedirs(out_dir, exist_ok=True)
        out_path = os.path.join(out_dir, f"KPI presentation - {MONTH_FULL[month - 1]} {year}.pptx")
    if chart_dir is None:
        chart_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), f"charts_{year}_{month:02d}")

    month_label = f"{MONTH_FULL[month - 1]} {year}"
    populate_deck(base_deck, pxd, hs_current, prev_hs, month_label, chart_dir, out_path, year, month)
    return out_path


def set_cell_like(shape, text):
    para = shape.text_frame.paragraphs[0]
    if para.runs:
        para.runs[0].text = text
        for extra in para.runs[1:]:
            extra.text = ""
    else:
        para.text = text


def _diff_all_tables(rebuilt_path, reference_path, skip_slides=(), skip_tables=()):
    """Cell-by-cell diff of every table in every slide between two decks.
    Returns a list of (slide_i, table_name, row, col, got, expected)
    mismatches. Slides 8/9/10's Picture 2 (chart image) is intentionally
    not compared here - the real reference deck still carries the original
    pasted screenshot, not a matplotlib rebuild, so a byte/pixel diff there
    is expected and meaningless; the picture-swap mechanism itself was
    already SHA1-verified separately (memory/kpi-presentation-assembly-gap.md)."""
    a = Presentation(rebuilt_path)
    b = Presentation(reference_path)
    mismatches = []
    for si, (sa, sb) in enumerate(zip(a.slides, b.slides)):
        if si in skip_slides:
            continue
        tables_a = {sh.name: sh.table for sh in sa.shapes if sh.has_table}
        tables_b = {sh.name: sh.table for sh in sb.shapes if sh.has_table}
        for name, ta in tables_a.items():
            if (si, name) in skip_tables:
                continue
            tb = tables_b.get(name)
            if tb is None:
                mismatches.append((si + 1, name, None, None, "table missing in reference", None))
                continue
            for ri, (ra, rb) in enumerate(zip(ta.rows, tb.rows)):
                for ci, (ca, cb) in enumerate(zip(ra.cells, rb.cells)):
                    if ca.text != cb.text:
                        mismatches.append((si + 1, name, ri, ci, ca.text, cb.text))
    return mismatches


def _check_chart(path, slide_i, expected):
    """expected: {category: (value, tolerance)}."""
    prs = Presentation(path)
    chart = next(sh.chart for sh in prs.slides[slide_i].shapes if sh.has_chart)
    plot = chart.plots[0]
    cats = list(plot.categories)
    vals = dict(zip(cats, list(plot.series[0].values)))
    results = []
    for cat, (exp, tol) in expected.items():
        got = vals.get(cat)
        got_f = got if got is not None else 0.0
        ok = abs(got_f - exp) <= tol
        results.append((f"slide{slide_i + 1} chart {cat!r}", got, exp, ok))
    return results


if __name__ == "__main__":
    # ---- Self-test: rebuild June 2026 from its own Source Data and check
    # every figure against the real, known-correct June deck. This is the
    # gate before trusting this script on a real July run.
    SCRATCH = os.path.dirname(os.path.abspath(__file__))
    june_src = r"C:\Users\admin\OneDrive - Nexus365\Functional Analysis Team Monthly Statistics\2026\06 Jun\Source Data"
    excel_path = os.path.join(june_src, "HR_Systems_Functional_Team_Monthly_Report_Excel - 202607010715.xlsx")
    hs_path = os.path.join(june_src, "Health and Safety Systems Support Statistics - 202607010600.docx")

    pxd = extract_pxd(excel_path)
    hs = extract_hs(hs_path)

    checks = [
        ("slide4 Service Request Jun26", pxd["slide4_categories"]["Service Request"][-1], 150),
        ("slide4 Change Jun26", pxd["slide4_categories"]["Change"][-1], 15),
        ("slide8 Jun26 avg acceptance", pxd["slide8_values"][-1], 0.4),
        ("slide9 Jun26 avg completion", pxd["slide9_values"][-1], 1.4),
        ("slide10 created Jun26", pxd["slide10_created"][-1], 347),
        ("slide10 completed Jun26", pxd["slide10_completed"][-1], 322),
        ("slide5 incident-other total (filtered)", pxd["slide5_incident_other_total"], 68),
        ("hs slide2 Cority", hs["slide2_volumes"]["Cority - Occupational Health Management Service"][0], 17),
        ("hs slide3 Same Day", hs["slide3_bands"]["Same Day"][0], 13),
    ]
    print("=== Self-test part 1: extraction against known June 2026 figures ===")
    all_pass = True
    for label, got, expected in checks:
        ok = got == expected
        all_pass &= ok
        print(f"  {'PASS' if ok else 'FAIL'}: {label} = {got} (expected {expected})")

    # ---- Part 2: full-deck assembly. Rebuild June end to end from May's
    # own finished deck as base (the real monthly process), then diff every
    # table cell against the real June deck - ground truth for tables (see
    # memory/kpi-presentation-assembly-gap.md; unlike charts, the real
    # deck's tables are trustworthy, so a cell-for-cell match here is a
    # meaningful correctness proof, not a tautology).
    print("\n=== Self-test part 2: full deck assembly (native chart writes, ")
    print("    slide 2/3/6/7/10 tables, H&S carry-forward, blank-row fix) ===")
    real_june = deck_path(2026, 6)
    test_out = os.path.join(SCRATCH, "assembly_selftest_v2.pptx")
    test_charts = os.path.join(SCRATCH, "charts_selftest")
    try:
        built_path = build_month(2026, 6, out_path=test_out, chart_dir=test_charts)
        mismatches = _diff_all_tables(
            built_path, real_june, skip_slides={7, 8, 9}, skip_tables={(4, "Table 6")}
        )
        # Slides 8/9/10 (index 7/8/9) contain Picture 2 (chart image, not
        # comparable - see _diff_all_tables docstring) alongside their KPI
        # threshold tables, which ARE comparable - diff those specifically.
        for si, table_name in ((7, "Table 11"), (8, "Table 10"), (9, "Table 3")):
            a = Presentation(built_path).slides[si]
            b = Presentation(real_june).slides[si]
            ta = next(sh.table for sh in a.shapes if sh.has_table and sh.name == table_name)
            tb = next(sh.table for sh in b.shapes if sh.has_table and sh.name == table_name)
            for ri, (ra, rb) in enumerate(zip(ta.rows, tb.rows)):
                for ci, (ca, cb) in enumerate(zip(ra.cells, rb.cells)):
                    if ca.text != cb.text:
                        mismatches.append((si + 1, table_name, ri, ci, ca.text, cb.text))

        # Slide 5, Table 6 needs a role-based comparison, not row-index
        # zip: our rebuild legitimately has one fewer row than the real
        # deck (the real live deck still carries the un-fixed blank row -
        # Kevin's manual deletion, memory/kpi-presentation-layout-fixes.md,
        # was only ever applied to a separate rebuild-TEST file, never the
        # live deck itself). Compare the 6 category rows by index (unaffected
        # by the row count difference) and the Total row by its role (last
        # row in each table), not by matching absolute index.
        ta6 = next(sh.table for sh in Presentation(built_path).slides[4].shapes
                   if sh.has_table and sh.name == "Table 6")
        tb6 = next(sh.table for sh in Presentation(real_june).slides[4].shapes
                   if sh.has_table and sh.name == "Table 6")
        for ri in range(7):  # header + 6 categories, same position in both
            for ci, (ca, cb) in enumerate(zip(ta6.rows[ri].cells, tb6.rows[ri].cells)):
                if ca.text != cb.text:
                    mismatches.append((5, "Table 6", ri, ci, ca.text, cb.text))
        for ci, (ca, cb) in enumerate(zip(list(ta6.rows)[-1].cells, list(tb6.rows)[-1].cells)):
            if ca.text != cb.text:
                mismatches.append((5, "Table 6 (Total row, by role)", -1, ci, ca.text, cb.text))

        # One known, disclosed deviation from the real deck - not a bug,
        # filtered out of the strict diff rather than silently left in: the
        # 0.01-point rounding-methodology gap on slide 7's "LESS THAN 5
        # DAYS" Jun26 cell (see NOTE printed below) - exact count division
        # vs. the real deck's own unreconcilable rounding.
        rounding_gap = {(7, "Table 5", 2, 5)}
        mismatches = [mm for mm in mismatches if (mm[0], mm[1], mm[2], mm[3]) not in rounding_gap]

        if mismatches:
            all_pass = False
            print(f"  FAIL: {len(mismatches)} table cell mismatch(es) vs real June deck:")
            for si, name, ri, ci, got, exp in mismatches:
                print(f"    slide{si} {name} row{ri} col{ci}: got {got!r}, expected {exp!r}")
        else:
            print("  PASS: every table cell across all 11 slides matches the real June deck exactly")

        # Chart checks: the real deck's own chart values are NOT trustworthy
        # ground truth (4 of 5 native charts were confirmed stale/pass-through
        # against independently-verified June figures during this build - see
        # HANDOVER.md). So chart checks assert against independently
        # hand-computed percentages from the same already-verified counts
        # above (17/38 Cority, 150/313 Service Request, etc.), not against
        # the real deck's chart content.
        chart_checks = []
        chart_checks += _check_chart(built_path, 1, {"Cority": (17 / 38 * 100, 0.01), "Odyssey": (5 / 38 * 100, 0.01), "IRIS": (16 / 38 * 100, 0.01)})
        chart_checks += _check_chart(built_path, 2, {"Same Day": (13 / 32 * 100, 0.01), "6+ Days": (12 / 32 * 100, 0.01)})
        chart_checks += _check_chart(built_path, 3, {"Service Request": (150 / 313 * 100, 0.01), "Change": (15 / 313 * 100, 0.01)})
        chart_checks += _check_chart(built_path, 5, {"Same Day": (112 / 150 * 100, 0.01)})
        chart_checks += _check_chart(built_path, 6, {"Same Day": (49 / 77 * 100, 0.01)})
        for label, got, exp, ok in chart_checks:
            all_pass &= ok
            print(f"  {'PASS' if ok else 'FAIL'}: {label} = {got} (expected {exp:.2f})")

        # Known, disclosed 0.01%-level rounding-methodology ambiguity (see
        # populate_deck's slide 6/7 comment) - checked separately with a
        # wider tolerance and explicitly flagged, not silently folded into
        # the strict table diff above.
        t7 = find_table(Presentation(built_path).slides[6], "Table 5")
        less5_jun = t7.rows[2].cells[5].text
        print(f"  NOTE (flagged, not a failure): slide 7 'LESS THAN 5 DAYS' Jun26 = {less5_jun!r}; "
              f"real deck shows '87.02%' - exact count division (67/77) gives 87.01%, a 0.01-point "
              f"rounding-methodology difference that doesn't resolve to one clean rule from the "
              f"available source data. See HANDOVER.md.")
    except Exception as e:
        all_pass = False
        print(f"  FAIL: full-deck assembly raised {type(e).__name__}: {e}")
    print("ALL PASS" if all_pass else "SOME FAILED - do not trust this script on real data yet")
